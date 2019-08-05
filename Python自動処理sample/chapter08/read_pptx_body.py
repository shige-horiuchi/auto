import pptx

prs = pptx.Presentation('book-proposal1.pptx')
for i, slide in enumerate(prs.slides):
    print(f'=' * 72)
    print(f'スライド: {i}')
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        for paragraph in shape.text_frame.paragraphs:
            print(paragraph.text)
